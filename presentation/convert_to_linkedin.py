#!/usr/bin/env python3
"""
Generate a LinkedIn carousel (portrait PPTX → export as PDF) for
"Snowflake AI for FinOps — Cortex in Practice" training.

10 slides, 4:5 portrait (1080×1350 px ≈ 7.50×9.38 in).
Uses the DevOpsDays Raleigh template for font embedding.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(SCRIPT_DIR, "DevOpsDaysRaleighWorkshopTemplate.pptx")
LOGO = os.path.join(SCRIPT_DIR, "devopsdays-logo.png")
OUTPUT = os.path.join(SCRIPT_DIR, "linkedin-carousel.pptx")

# ── Slide dimensions: 4:5 portrait ──
SLIDE_W = Emu(6858000)   # 7.50 in
SLIDE_H = Emu(8572500)   # 9.375 in  (7.50 × 1.25)

# ── Colors (DevOpsDays palette) ──
BLUE = RGBColor(0x05, 0x8D, 0xC7)
GREEN = RGBColor(0x50, 0xB4, 0x32)
ORANGE = RGBColor(0xED, 0x56, 0x1B)
CYAN = RGBColor(0x24, 0xCB, 0xE5)
TEAL = RGBColor(0x15, 0x81, 0x58)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_BG = RGBColor(0xF3, 0xF3, 0xF3)
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
DANGER = RGBColor(0xCC, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF4, 0xFD)

# ── Fonts ──
FONT_HEADING = "Montserrat"
FONT_BODY = "Montserrat"
FONT_CODE = "Courier New"

# ── Layout constants for portrait slides ──
MARGIN = Inches(0.5)
CONTENT_W = Inches(6.5)  # 7.5 - 2×0.5


# ════════════════════════════════════════════════════════════════
#  Helpers
# ════════════════════════════════════════════════════════════════

def delete_existing_slides(prs):
    """Remove all existing slides from the template."""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn("r:id"))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=TEXT_COLOR, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY, word_wrap=True):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height):
    """Add empty text box, return (shape, text_frame)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def set_shape_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_shape_border(shape, color, width_pt=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def add_pill(slide, left, top, width, height, text,
             bg_color=BLUE, text_color=WHITE, font_size=11):
    """Add a rounded-rect pill/badge."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    set_shape_fill(shape, bg_color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = FONT_HEADING
    return shape


def add_card(slide, left, top, width, height, title, bullets,
             title_color=BLUE, border_color=CARD_BORDER, border_width=1):
    """Rounded-rect card with title + bullet list."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    set_shape_fill(shape, GRAY_BG)
    set_shape_border(shape, border_color, border_width)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.05)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = title_color
    run.font.name = FONT_HEADING
    p.space_after = Pt(4)

    for bullet in bullets:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"\u25b8 {bullet}"
        run.font.size = Pt(11)
        run.font.color.rgb = TEXT_COLOR
        run.font.name = FONT_BODY
        p.space_before = Pt(2)
        p.space_after = Pt(2)

    return shape


def add_table(slide, data, left, top, width, height, col_widths=None):
    """Add a table from a 2-D list. First row is header."""
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT_BODY
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = TEXT_COLOR

        # Header row fill
        if r == 0:
            for c in range(cols):
                cell = table.cell(0, c)
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn("a:solidFill"), {})
                srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": "058DC7"})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

    return tbl_shape


def add_slide_number(slide, num, total=10):
    """Small slide counter at bottom."""
    add_textbox(slide, Inches(3.0), Inches(8.9), Inches(1.5), Inches(0.3),
                f"{num} / {total}", font_size=9, color=MUTED,
                alignment=PP_ALIGN.CENTER)


def add_divider_line(slide, top, color=BLUE):
    """Horizontal accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN, top, CONTENT_W, Inches(0.04)
    )
    set_shape_fill(shape, color)
    shape.line.fill.background()


# ════════════════════════════════════════════════════════════════
#  Slide builders
# ════════════════════════════════════════════════════════════════

def build_slide_01_hook(prs):
    """SLIDE 1: Title / Hook — bold opening statement."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # BLANK

    # Blue background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

    # Top accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_W, Inches(0.15)
    )
    set_shape_fill(bar, BLUE)
    bar.line.fill.background()

    # Logo
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(2.75), Inches(1.0), Inches(2.0))

    # Event tag
    add_textbox(slide, MARGIN, Inches(2.2), CONTENT_W, Inches(0.35),
                "DEVOPSDAYS RALEIGH \u2014 APR 30 \u2013 MAY 1, 2026",
                font_size=11, bold=True, color=ORANGE,
                alignment=PP_ALIGN.CENTER)

    # Hook line
    add_textbox(slide, MARGIN, Inches(3.2), CONTENT_W, Inches(1.0),
                "AI Is Now\na Cost Center.",
                font_size=36, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Training title
    add_textbox(slide, MARGIN, Inches(4.8), CONTENT_W, Inches(0.6),
                "Snowflake AI for FinOps",
                font_size=22, bold=True, color=BLUE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN, Inches(5.4), CONTENT_W, Inches(0.4),
                "Cortex in Practice",
                font_size=16, color=CYAN,
                alignment=PP_ALIGN.CENTER)

    # Metadata
    add_textbox(slide, MARGIN, Inches(6.3), CONTENT_W, Inches(0.35),
                "90 Minutes  |  9 Modules  |  Hands-On Labs",
                font_size=11, color=MUTED,
                alignment=PP_ALIGN.CENTER)

    # Authors
    add_textbox(slide, MARGIN, Inches(7.0), CONTENT_W, Inches(0.35),
                "Velu Natarajan & Krishnakumar Mohanram",
                font_size=11, color=RGBColor(0xAA, 0xAA, 0xAA),
                alignment=PP_ALIGN.CENTER)

    # Swipe prompt
    add_textbox(slide, MARGIN, Inches(8.4), CONTENT_W, Inches(0.3),
                "Swipe \u25B6 to see what we covered",
                font_size=10, color=MUTED,
                alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 1)


def build_slide_02_problem(prs):
    """SLIDE 2: The Problem — why AI costs are different."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    add_textbox(slide, MARGIN, Inches(0.4), CONTENT_W, Inches(0.5),
                "The Problem", font_size=28, bold=True, color=BLUE,
                alignment=PP_ALIGN.LEFT)

    add_divider_line(slide, Inches(1.0))

    add_textbox(slide, MARGIN, Inches(1.3), CONTENT_W, Inches(0.8),
                "AI spend is invisible to traditional\nSnowflake monitoring.",
                font_size=18, bold=True, color=DARK,
                alignment=PP_ALIGN.LEFT)

    problems = [
        ("Credit Burn", DANGER,
         "AI credit consumption scales faster than traditional compute. "
         "No warehouse to suspend."),
        ("Behavior", ORANGE,
         "Teams switch models frequently. Usage patterns shift weekly. "
         "Exploratory prompts run at scale."),
        ("Visibility", BLUE,
         "Serverless AI spend is invisible to warehouse dashboards. "
         "Cortex AI runs outside warehouses."),
        ("Signals", GREEN,
         "Same prompt on a large model costs 10\u00d7 more. "
         "Without tracking model choice, costs stay hidden."),
    ]

    y = Inches(2.5)
    for name, color, desc in problems:
        # Color bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN, y, Inches(0.12), Inches(1.1)
        )
        set_shape_fill(bar, color)
        bar.line.fill.background()

        add_textbox(slide, Inches(0.8), y, Inches(5.5), Inches(0.35),
                    name, font_size=14, bold=True, color=color)
        add_textbox(slide, Inches(0.8), y + Inches(0.35), Inches(5.5), Inches(0.7),
                    desc, font_size=10, color=TEXT_COLOR)

        y += Inches(1.4)

    # Bottom callout
    txBox, tf = add_rich_textbox(slide, MARGIN, Inches(8.1), CONTENT_W, Inches(0.6))
    set_shape_fill(txBox, RGBColor(0xFD, 0xED, 0xED))
    set_shape_border(txBox, DANGER, 1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "A misconfigured COMPLETE() loop across 1M rows can exceed a week of warehouse compute in minutes."
    run.font.size = Pt(10)
    run.font.color.rgb = DANGER
    run.font.name = FONT_BODY
    run.font.bold = True

    add_slide_number(slide, 2)


def build_slide_03_two_layers(prs):
    """SLIDE 3: Two Cost Layers — warehouse vs AI tokens."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    add_textbox(slide, MARGIN, Inches(0.4), CONTENT_W, Inches(0.5),
                "Two Cost Layers", font_size=28, bold=True, color=BLUE)

    add_divider_line(slide, Inches(1.0))

    add_textbox(slide, MARGIN, Inches(1.3), CONTENT_W, Inches(0.5),
                "Snowflake now has TWO independent billing layers.",
                font_size=14, color=TEXT_COLOR)

    # Layer 1 card
    add_card(slide, MARGIN, Inches(2.2), CONTENT_W, Inches(2.5),
             "Layer 1: Warehouse Compute",
             ["Traditional Snowflake billing",
              "Based on warehouse size \u00d7 time",
              "Stops when warehouse suspends",
              "Governed by Resource Monitors"],
             title_color=BLUE, border_color=BLUE, border_width=2)

    # Layer 2 card
    add_card(slide, MARGIN, Inches(5.0), CONTENT_W, Inches(2.5),
             "Layer 2: AI Token Credits",
             ["New billing layer for Cortex AI",
              "Based on tokens consumed per call",
              "Does NOT stop on suspend",
              "Governed by Budgets + Credit Limits"],
             title_color=DANGER, border_color=DANGER, border_width=2)

    # Key insight
    txBox, tf = add_rich_textbox(slide, MARGIN, Inches(7.8), CONTENT_W, Inches(0.7))
    set_shape_fill(txBox, LIGHT_BLUE_BG)
    set_shape_border(txBox, BLUE, 2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Resource Monitors only cover Layer 1.\nAI token credits need separate governance."
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = DARK
    run.font.name = FONT_BODY

    add_slide_number(slide, 3)


def build_slide_04_what_we_built(prs):
    """SLIDE 4: What We Built — 9-module training overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    add_textbox(slide, MARGIN, Inches(0.4), CONTENT_W, Inches(0.5),
                "What We Built", font_size=28, bold=True, color=BLUE)

    add_divider_line(slide, Inches(1.0))

    add_textbox(slide, MARGIN, Inches(1.2), CONTENT_W, Inches(0.4),
                "Understand \u2192 Use \u2192 Track \u2192 Govern AI Spend",
                font_size=13, bold=True, color=TEAL,
                alignment=PP_ALIGN.CENTER)

    modules = [
        ("01", "AI Token Economy", "Concept", BLUE),
        ("02", "Cortex AI Capabilities", "Concept", BLUE),
        ("03", "Cortex Code Setup", "Setup", ORANGE),
        ("04", "Environment Setup", "Lab", GREEN),
        ("05", "AI SQL Hands-On", "Lab", GREEN),
        ("06", "Usage Tracking", "Lab", GREEN),
        ("07", "AI Credits Transition", "Lab", GREEN),
        ("08", "Streamlit Dashboard", "Lab", GREEN),
        ("09", "Closing & Governance", "Wrap-up", MUTED),
    ]

    y = Inches(1.8)
    for num, name, typ, color in modules:
        # Number pill
        add_pill(slide, MARGIN, y, Inches(0.55), Inches(0.5),
                 num, bg_color=color, font_size=12)

        # Module name
        add_textbox(slide, Inches(1.2), y, Inches(4.0), Inches(0.5),
                    name, font_size=13, bold=True, color=DARK)

        # Type badge
        add_pill(slide, Inches(5.4), y + Inches(0.05), Inches(1.3), Inches(0.4),
                 typ, bg_color=color, font_size=9)

        y += Inches(0.6)

    # Bottom summary
    add_textbox(slide, MARGIN, Inches(8.0), CONTENT_W, Inches(0.5),
                "90 minutes \u2022 5 hands-on labs \u2022 13 SQL queries \u2022 1 Streamlit dashboard",
                font_size=10, color=MUTED,
                alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 4)


def build_slide_05_shadow_waste(prs):
    """SLIDE 5: Shadow Waste — the #1 insight."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    add_textbox(slide, MARGIN, Inches(0.4), CONTENT_W, Inches(0.5),
                "Shadow Waste", font_size=28, bold=True, color=DANGER)

    add_divider_line(slide, Inches(1.0), color=DANGER)

    add_textbox(slide, MARGIN, Inches(1.3), CONTENT_W, Inches(0.8),
                "Using COMPLETE() for tasks with\nspecialized functions is the #1\nsource of hidden AI spend.",
                font_size=16, bold=True, color=DARK)

    # Savings table
    savings_data = [
        ["Task", "COMPLETE()", "Specialized", "Savings"],
        ["Sentiment", "0.30 cr", "0.02 cr", "93%"],
        ["Translate", "0.45 cr", "0.05 cr", "89%"],
        ["Classify", "0.30 cr", "0.03 cr", "90%"],
        ["Extract", "0.30 cr", "0.04 cr", "87%"],
    ]

    add_table(slide, savings_data,
              MARGIN, Inches(2.5), CONTENT_W, Inches(2.5),
              col_widths=[Inches(1.5), Inches(1.7), Inches(1.7), Inches(1.6)])

    # 8 detection patterns
    add_textbox(slide, MARGIN, Inches(5.3), CONTENT_W, Inches(0.4),
                "8 Detection Patterns in Module 06:",
                font_size=13, bold=True, color=BLUE)

    patterns = [
        "Over-sized models (40\u201375% savings potential)",
        "Redundant / duplicate AI calls",
        "Prompt bloat (rising avg tokens/call)",
        "Agent loops (MAX vs AVG credit gap)",
        "Dev-in-prod (LATERAL FLATTEN role detection)",
        "Idle Search indexing (100% waste)",
        "WoW anomaly spikes (LAG-7 detection)",
        "Untagged spend (attribution coverage gaps)",
    ]

    y = Inches(5.8)
    for pat in patterns:
        add_textbox(slide, Inches(0.7), y, Inches(6.0), Inches(0.3),
                    f"\u25b8 {pat}", font_size=9, color=TEXT_COLOR)
        y += Inches(0.3)

    add_slide_number(slide, 5)


def build_slide_06_seven_views(prs):
    """SLIDE 6: 7 ACCOUNT_USAGE Views — the monitoring toolkit."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    add_textbox(slide, MARGIN, Inches(0.4), CONTENT_W, Inches(0.5),
                "7 ACCOUNT_USAGE Views", font_size=28, bold=True, color=BLUE)

    add_divider_line(slide, Inches(1.0))

    add_textbox(slide, MARGIN, Inches(1.2), CONTENT_W, Inches(0.4),
                "Your complete AI cost visibility toolkit",
                font_size=13, color=MUTED)

    views_data = [
        ["View", "Tracks"],
        ["CORTEX_AI_FUNCTIONS\n_USAGE_HISTORY", "LLM / task functions"],
        ["CORTEX_AISQL\n_USAGE_HISTORY", "AI SQL tokens"],
        ["DOCUMENT_AI\n_USAGE_HISTORY", "Document AI"],
        ["CORTEX_SEARCH_DAILY\n_USAGE_HISTORY", "Search services"],
        ["METERING_DAILY\n_HISTORY", "All SERVICE_TYPEs"],
        ["CORTEX_CODE_CLI\n_USAGE_HISTORY", "Cortex Code CLI"],
        ["CORTEX_CODE_SNOWSIGHT\n_USAGE_HISTORY", "Cortex Code Snowsight"],
    ]

    add_table(slide, views_data,
              MARGIN, Inches(1.8), CONTENT_W, Inches(5.2),
              col_widths=[Inches(4.0), Inches(2.5)])

    # Tip
    txBox, tf = add_rich_textbox(slide, MARGIN, Inches(7.3), CONTENT_W, Inches(0.9))
    set_shape_fill(txBox, RGBColor(0xFD, 0xED, 0xED))
    set_shape_border(txBox, DANGER, 1)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "USER_ID in Cortex views is a NUMBER.\nJOIN to SNOWFLAKE.ACCOUNT_USAGE.USERS for names.\n~45 min latency on all views."
    run.font.size = Pt(10)
    run.font.color.rgb = DANGER
    run.font.name = FONT_BODY

    add_slide_number(slide, 6)


def build_slide_07_guardrails(prs):
    """SLIDE 7: Credit Guardrails — three-layer governance."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    add_textbox(slide, MARGIN, Inches(0.4), CONTENT_W, Inches(0.5),
                "Credit Guardrails", font_size=28, bold=True, color=BLUE)

    add_divider_line(slide, Inches(1.0))

    add_textbox(slide, MARGIN, Inches(1.2), CONTENT_W, Inches(0.4),
                "Three layers of Cortex Code governance",
                font_size=13, color=MUTED)

    # Layer cards
    layers = [
        ("PROACTIVE", "Per-User Credit Limits",
         "Hard block \u2014 rolling 24h window\nAccount-wide default + per-user override",
         DANGER),
        ("THRESHOLD", "Monthly Budgets",
         "Alert before overshoot \u2014 no hard block\nEarly warning system for all spend",
         ORANGE),
        ("AUDIT", "Usage History Views",
         "Attribution by user + model after the fact\nCLI and Snowsight tracked separately",
         BLUE),
    ]

    y = Inches(1.8)
    for label, title, desc, color in layers:
        # Label pill
        add_pill(slide, MARGIN, y, Inches(1.8), Inches(0.4),
                 label, bg_color=color, font_size=10)

        # Card
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MARGIN, y + Inches(0.5), CONTENT_W, Inches(1.3)
        )
        set_shape_fill(shape, GRAY_BG)
        set_shape_border(shape, color, 2)
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = FONT_HEADING
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(10)
        run2.font.color.rgb = TEXT_COLOR
        run2.font.name = FONT_BODY

        y += Inches(2.0)

    # Parameters footnote
    add_textbox(slide, MARGIN, Inches(8.0), CONTENT_W, Inches(0.7),
                "Parameters:\n"
                "CORTEX_CODE_CLI_DAILY_EST_CREDIT_LIMIT_PER_USER\n"
                "CORTEX_CODE_SNOWSIGHT_DAILY_EST_CREDIT_LIMIT_PER_USER",
                font_size=8, color=MUTED)

    add_slide_number(slide, 7)


def build_slide_08_habits(prs):
    """SLIDE 8: Three FinOps Habits — the takeaway framework."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    add_textbox(slide, MARGIN, Inches(0.4), CONTENT_W, Inches(0.5),
                "Three FinOps Habits", font_size=28, bold=True, color=BLUE)

    add_divider_line(slide, Inches(1.0))

    add_textbox(slide, MARGIN, Inches(1.2), CONTENT_W, Inches(0.4),
                "The framework for governing AI spend",
                font_size=13, color=MUTED)

    habits = [
        ("ESTIMATE", BLUE,
         "Test on 10 rows.\nProject to full dataset.\nDecide.",
         "Never run AI at scale without first testing on a small sample, "
         "measuring the cost, and projecting."),
        ("MONITOR", GREEN,
         "Dashboard is your\nearly warning system.",
         "The Streamlit dashboard you built should be checked weekly. "
         "Set alert thresholds before problems become emergencies."),
        ("CHOOSE", ORANGE,
         "The cheapest model that\nmeets your quality bar.",
         "Model selection is a cost decision. Upgrade only when "
         "the output genuinely requires higher capability."),
    ]

    y = Inches(1.8)
    for name, color, headline, desc in habits:
        # Big habit card
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MARGIN, y, CONTENT_W, Inches(2.0)
        )
        set_shape_fill(shape, GRAY_BG)
        set_shape_border(shape, color, 3)
        shape.shadow.inherit = False

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.12)

        # Name
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = name
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = FONT_HEADING
        p.space_after = Pt(4)

        # Headline
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = headline
        run2.font.size = Pt(12)
        run2.font.bold = True
        run2.font.color.rgb = DARK
        run2.font.name = FONT_BODY
        p2.space_after = Pt(4)

        # Description
        p3 = tf.add_paragraph()
        run3 = p3.add_run()
        run3.text = desc
        run3.font.size = Pt(9)
        run3.font.color.rgb = MUTED
        run3.font.name = FONT_BODY

        y += Inches(2.2)

    add_slide_number(slide, 8)


def build_slide_09_quote(prs):
    """SLIDE 9: Key stat / closing quote."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

    # Top accent
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_W, Inches(0.15)
    )
    set_shape_fill(bar, BLUE)
    bar.line.fill.background()

    # Big quote
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.9), Inches(2.5),
                "AI token consumption\nis the fastest-growing\nline item in cloud\ndata platform bills\nin 2025.",
                font_size=24, bold=True, color=WHITE,
                alignment=PP_ALIGN.LEFT)

    # Divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(5.0), Inches(2.0), Inches(0.06)
    )
    set_shape_fill(divider, BLUE)
    divider.line.fill.background()

    # Sub-quote
    add_textbox(slide, Inches(0.8), Inches(5.3), Inches(5.9), Inches(1.0),
                "The teams who govern AI costs now\nwill be the ones who can afford\nto scale AI later.",
                font_size=14, color=CYAN)

    # Key numbers
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(2.5), Inches(0.5),
                "$3/credit", font_size=18, bold=True, color=ORANGE)
    add_textbox(slide, Inches(3.5), Inches(7.0), Inches(3.0), Inches(0.5),
                "87\u201393% savings possible",
                font_size=18, bold=True, color=GREEN)

    add_textbox(slide, Inches(0.8), Inches(7.6), Inches(5.9), Inches(0.4),
                "by choosing specialized functions over COMPLETE()",
                font_size=10, color=MUTED)

    add_slide_number(slide, 9)


def build_slide_10_cta(prs):
    """SLIDE 10: Call to action + credits."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Blue background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLUE

    # Logo
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(2.75), Inches(0.8), Inches(2.0))

    # Event
    add_textbox(slide, MARGIN, Inches(2.2), CONTENT_W, Inches(0.4),
                "DEVOPSDAYS RALEIGH",
                font_size=14, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN, Inches(2.6), CONTENT_W, Inches(0.4),
                "April 30 \u2013 May 1, 2026",
                font_size=12, color=RGBColor(0xB0, 0xE0, 0xFF),
                alignment=PP_ALIGN.CENTER)

    # Training title
    add_textbox(slide, MARGIN, Inches(3.5), CONTENT_W, Inches(0.6),
                "Snowflake AI for FinOps",
                font_size=26, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN, Inches(4.1), CONTENT_W, Inches(0.4),
                "Cortex in Practice",
                font_size=16, color=CYAN,
                alignment=PP_ALIGN.CENTER)

    # What's included
    items = [
        "9 modules \u2022 90 minutes \u2022 hands-on labs",
        "13 SQL queries \u2022 Streamlit FinOps dashboard",
        "7 ACCOUNT_USAGE views \u2022 credit guardrails",
        "Shadow waste detection \u2022 cost governance",
    ]

    y = Inches(5.2)
    for item in items:
        add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.35),
                    item, font_size=11, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        y += Inches(0.35)

    # Authors
    add_textbox(slide, MARGIN, Inches(7.2), CONTENT_W, Inches(0.35),
                "Velu Natarajan & Krishnakumar Mohanram",
                font_size=12, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # CTA
    cta_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(8.0), Inches(4.5), Inches(0.55)
    )
    set_shape_fill(cta_shape, ORANGE)
    cta_shape.line.fill.background()
    cta_shape.shadow.inherit = False
    tf = cta_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Like \u2764 + Comment to get the materials"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_HEADING

    add_slide_number(slide, 10)


# ════════════════════════════════════════════════════════════════
#  Main
# ════════════════════════════════════════════════════════════════

def main():
    print("Loading template:", TEMPLATE)
    prs = Presentation(TEMPLATE)

    # Override slide dimensions to portrait 4:5
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Removing existing template slides...")
    delete_existing_slides(prs)

    print("Building 10 LinkedIn carousel slides...")

    builders = [
        (build_slide_01_hook, "Hook / Title"),
        (build_slide_02_problem, "The Problem"),
        (build_slide_03_two_layers, "Two Cost Layers"),
        (build_slide_04_what_we_built, "What We Built"),
        (build_slide_05_shadow_waste, "Shadow Waste"),
        (build_slide_06_seven_views, "7 ACCOUNT_USAGE Views"),
        (build_slide_07_guardrails, "Credit Guardrails"),
        (build_slide_08_habits, "Three FinOps Habits"),
        (build_slide_09_quote, "Key Quote"),
        (build_slide_10_cta, "CTA / Credits"),
    ]

    for i, (builder, label) in enumerate(builders, 1):
        builder(prs)
        print(f"  Slide {i:2d}: {label}")

    print(f"\nSaving to: {OUTPUT}")
    prs.save(OUTPUT)
    print(f"Done! {len(prs.slides)} slides, portrait 4:5 format.")
    print("\nNext steps:")
    print("  1. Open in PowerPoint / Google Slides")
    print("  2. File > Export as PDF")
    print("  3. Upload PDF to LinkedIn as a Document post")


if __name__ == "__main__":
    main()
