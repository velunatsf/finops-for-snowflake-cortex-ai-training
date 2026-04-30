#!/usr/bin/env python3
"""
Convert DevOpsDays HTML presentation to PPTX using the DevOpsDays Raleigh template.
Uses python-pptx to programmatically build 14 slides matching the HTML content.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(SCRIPT_DIR, "DevOpsDaysRaleighWorkshopTemplate.pptx")
LOGO = os.path.join(SCRIPT_DIR, "devopsdays-logo.png")
GEAR = os.path.join(SCRIPT_DIR, "gear-watermark.png")
OUTPUT = os.path.join(SCRIPT_DIR, "presentation-devopsdays.pptx")

# ── Colors (from PPTX theme) ──
BLUE = RGBColor(0x05, 0x8D, 0xC7)
GREEN = RGBColor(0x50, 0xB4, 0x32)
ORANGE = RGBColor(0xED, 0x56, 0x1B)
CYAN = RGBColor(0x24, 0xCB, 0xE5)
TEAL = RGBColor(0x15, 0x81, 0x58)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_BG = RGBColor(0xF3, 0xF3, 0xF3)
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
DANGER = RGBColor(0xCC, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
BLACK = RGBColor(0x00, 0x00, 0x00)

# ── Slide dimensions (from template) ──
SLIDE_W = Emu(12192000)  # 13.33"
SLIDE_H = Emu(6858000)   # 7.50"

# ── Fonts ──
FONT_HEADING = "Montserrat"
FONT_BODY = "Montserrat"
FONT_CODE = "Courier New"


# ════════════════════════════════════════════════════════════════
#  Helper functions
# ════════════════════════════════════════════════════════════════

def delete_existing_slides(prs):
    """Remove all existing slides from the template."""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn("r:id"))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=TEXT_COLOR, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY, word_wrap=True):
    """Add a simple text box and return the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height):
    """Add an empty text box and return the text frame for manual paragraph building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def set_shape_fill(shape, color):
    """Set solid fill on a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_shape_border(shape, color, width_pt=1):
    """Set border on a shape."""
    ln = shape.line
    ln.color.rgb = color
    ln.width = Pt(width_pt)


def add_card(slide, left, top, width, height, title, bullets, title_color=BLUE,
             border_color=CARD_BORDER, border_width=1, accent_side=None, accent_color=None):
    """Add a rounded-rect card with title and bullet list."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    set_shape_fill(shape, GRAY_BG)
    set_shape_border(shape, border_color, border_width)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.05)

    # Title paragraph
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = title_color
    run.font.name = FONT_HEADING
    p.space_after = Pt(4)

    # Bullet paragraphs
    for bullet in bullets:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"\u25b8 {bullet}"
        run.font.size = Pt(10)
        run.font.color.rgb = TEXT_COLOR
        run.font.name = FONT_BODY
        p.space_before = Pt(1)
        p.space_after = Pt(1)

    return shape


def add_slide_title(slide, text, font_size=28, color=BLUE):
    """Set the title placeholder text on a slide that has one."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            for p in ph.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.bold = True
                    run.font.color.rgb = color
                    run.font.name = FONT_HEADING
            return ph
    return None


def add_subtitle(slide, text, font_size=14, color=MUTED):
    """Add a subtitle line below the title area."""
    txBox = add_textbox(slide, Inches(0.7), Inches(1.1), Inches(11.5), Inches(0.4),
                        text, font_size=font_size, color=color)
    return txBox


def add_footer(slide):
    """Add #DevOpsDays footer text on bottom-right."""
    add_textbox(slide, Inches(10.5), Inches(7.0), Inches(2.5), Inches(0.35),
                "#DevOpsDays", font_size=10, bold=True, color=MUTED,
                alignment=PP_ALIGN.RIGHT)


def add_warning_box(slide, text, left, top, width, height=Inches(0.5)):
    """Add a red-bordered warning callout."""
    txBox, tf = add_rich_textbox(slide, left, top, width, height)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = DANGER
    run.font.name = FONT_BODY
    # Red left border via shape line
    set_shape_fill(txBox, RGBColor(0xFD, 0xED, 0xED))
    set_shape_border(txBox, DANGER, 1)
    return txBox


def add_highlight_box(slide, text, left, top, width, height=Inches(0.45)):
    """Add a blue highlight callout box."""
    txBox, tf = add_rich_textbox(slide, left, top, width, height)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.color.rgb = DARK
    run.font.name = FONT_BODY
    run.font.bold = True
    set_shape_fill(txBox, RGBColor(0xE8, 0xF4, 0xFD))
    set_shape_border(txBox, BLUE, 2)
    return txBox


def add_code_box(slide, code_text, left, top, width, height):
    """Add a dark code block text box."""
    txBox, tf = add_rich_textbox(slide, left, top, width, height)
    set_shape_fill(txBox, DARK)
    set_shape_border(txBox, RGBColor(0x44, 0x44, 0x66), 1)
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)

    for i, line in enumerate(code_text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        run.font.name = FONT_CODE
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    return txBox


def add_table(slide, rows_data, left, top, width, height, col_widths=None,
              header_color=BLUE, header_bg=None):
    """Add a table to the slide. rows_data[0] = header row."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths if provided
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT_BODY
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = header_color
                else:
                    p.font.color.rgb = TEXT_COLOR
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)

    return table_shape


def add_badge_text(paragraph, text, badge_color, font_size=9):
    """Add a colored badge-style run to an existing paragraph."""
    run = paragraph.add_run()
    run.text = f"  [{text}]"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = badge_color
    run.font.name = FONT_BODY


# ════════════════════════════════════════════════════════════════
#  Build slides
# ════════════════════════════════════════════════════════════════

def build_slide_01_title(prs):
    """SLIDE 1: Title slide with logo."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # TITLE layout

    # Logo
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(5.7), Inches(0.3), Inches(1.9))

    # Title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Snowflake AI for FinOps"
            for p in ph.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(40)
                    run.font.bold = True
                    run.font.color.rgb = DARK
                    run.font.name = FONT_HEADING
        elif ph.placeholder_format.idx == 1:
            ph.text = "Cortex in Practice"
            for p in ph.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(24)
                    run.font.bold = True
                    run.font.color.rgb = BLUE
                    run.font.name = FONT_HEADING

    # Metadata row
    add_textbox(slide, Inches(2.5), Inches(5.0), Inches(8.0), Inches(0.4),
                "90 Minutes    |    9 Modules    |    Hands-On Labs",
                font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Event tag
    add_textbox(slide, Inches(2.5), Inches(5.5), Inches(8.0), Inches(0.35),
                "DEVOPSDAYS RALEIGH \u2014 APR 30 - MAY 1, 2026",
                font_size=12, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

    # Authors
    add_textbox(slide, Inches(2.5), Inches(5.95), Inches(8.0), Inches(0.35),
                "Velu Natarajan & Krishnakumar Mohanram",
                font_size=12, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)

    add_footer(slide)


def build_slide_02_agenda(prs):
    """SLIDE 2: Training Agenda — 9-row table."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])  # TITLE_ONLY
    add_slide_title(slide, "Training Agenda")

    agenda_data = [
        ["#", "Module", "Duration", "Type"],
        ["01", "AI Token Economy", "10 min", "Concept"],
        ["02", "Cortex AI Capabilities", "15 min", "Concept"],
        ["03", "Cortex Code Setup", "5 min", "Setup"],
        ["04", "Environment Setup", "10 min", "Lab"],
        ["05", "AI SQL \u2014 Hands-On Exercises", "20 min", "Lab"],
        ["06", "Token & Credit Usage Tracking", "12 min", "Lab"],
        ["07", "AI Credits Transition", "5 min", "Lab"],
        ["08", "Streamlit FinOps Dashboard", "13 min", "Lab"],
        ["09", "Closing Note", "5 min", "Wrap-up"],
    ]

    tbl_shape = add_table(slide, agenda_data,
                          Inches(1.0), Inches(1.5), Inches(11.0), Inches(5.0),
                          col_widths=[Inches(0.8), Inches(6.5), Inches(1.5), Inches(2.2)])

    # Color the type badges
    type_colors = {
        "Concept": BLUE, "Setup": ORANGE, "Lab": GREEN, "Wrap-up": MUTED
    }
    table = tbl_shape.table
    for r in range(1, len(agenda_data)):
        # Module number in blue
        for p in table.cell(r, 0).text_frame.paragraphs:
            p.font.color.rgb = BLUE
            p.font.bold = True
        # Duration in muted
        for p in table.cell(r, 3).text_frame.paragraphs:
            type_text = agenda_data[r][3]
            p.font.color.rgb = type_colors.get(type_text, MUTED)
            p.font.bold = True

    add_footer(slide)


def build_slide_03_cost_layers(prs):
    """SLIDE 3: AI Is Now a Cost Center — Two cost layers."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])  # TITLE_ONLY
    add_slide_title(slide, "AI Is Now a Cost Center")

    # Left card — Warehouse Compute
    add_card(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(2.8),
             "Layer 1: Warehouse Compute",
             ["Traditional Snowflake billing",
              "Based on warehouse size \u00d7 time",
              "Stops when warehouse suspends"],
             title_color=BLUE, border_color=BLUE, border_width=2)

    # Right card — AI Token Credits
    add_card(slide, Inches(6.7), Inches(1.5), Inches(5.5), Inches(2.8),
             "Layer 2: AI Token Credits",
             ["New billing layer for Cortex AI",
              "Based on tokens consumed per call",
              "Does NOT stop on suspend"],
             title_color=DANGER, border_color=DANGER, border_width=2)

    # Warning box
    add_warning_box(slide,
                    "A misconfigured COMPLETE() loop across 1M rows can exceed "
                    "a week of warehouse compute in minutes.",
                    Inches(0.7), Inches(4.7), Inches(11.5), Inches(0.55))

    add_footer(slide)


def build_slide_04_capability_map(prs):
    """SLIDE 4: Cortex AI Capability Map — 2x2 grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])  # TITLE_ONLY
    add_slide_title(slide, "Cortex AI Capability Map")
    add_subtitle(slide, "Module 02 \u2014 Four categories of AI in Snowflake")

    cards = [
        (Inches(0.7), Inches(1.6), "AI Functions",
         ["AI_COMPLETE() \u2014 General LLM",
          "AI_SENTIMENT() \u2014 Lowest cost",
          "AI_SUMMARIZE_AGG() \u00b7 AI_TRANSLATE()",
          "AI_CLASSIFY() \u00b7 AI_EXTRACT()"],
         BLUE, CARD_BORDER),
        (Inches(6.7), Inches(1.6), "Search & Embed",
         ["Cortex Search (standing cost)",
          "AI_EMBED()",
          "Cortex Analyst"],
         BLUE, CARD_BORDER),
        (Inches(0.7), Inches(4.0), "ML Functions",
         ["FORECAST() \u00b7 ANOMALY_DETECTION()",
          "CLASSIFICATION()",
          "CONTRIBUTION_EXPLORER()"],
         BLUE, CARD_BORDER),
        (Inches(6.7), Inches(4.0), "Cortex Code",
         ["Natural language \u2192 SQL generation",
          "CLI + Snowsight \u00b7 Agentic workflows",
          "Own SERVICE_TYPEs for cost tracking"],
         BLUE, BLUE),
    ]

    for left, top, title, bullets, title_clr, border_clr in cards:
        add_card(slide, left, top, Inches(5.5), Inches(2.1),
                 title, bullets, title_color=title_clr,
                 border_color=border_clr, border_width=2 if border_clr == BLUE and title == "Cortex Code" else 1)

    # Strategy note
    add_textbox(slide, Inches(0.7), Inches(6.3), Inches(11.5), Inches(0.4),
                "Strategy: Start with specialized functions before reaching for COMPLETE(). "
                "They are optimized and cheaper.",
                font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

    add_footer(slide)


def build_slide_05_model_tiers(prs):
    """SLIDE 5: Model Cost Tiers — Budget / Standard / Premium."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    add_slide_title(slide, "Model Cost Tiers")
    add_subtitle(slide, "Same task, different models = 10-50x cost difference")

    col_w = Inches(3.7)
    col_h = Inches(3.0)
    top = Inches(1.6)
    gap = Inches(0.2)
    left_start = Inches(0.7)

    tiers = [
        ("Budget", GREEN, [
            "mistral-7b \u2014 8.3M tok/cr",
            "llama3-8b \u2014 5.3M tok/cr",
            "gemma-7b \u2014 8.3M tok/cr",
        ], "Extraction, classification, high volume"),
        ("Standard", BLUE, [
            "llama3-70b \u2014 826K tok/cr",
            "mixtral-8x7b \u2014 4.5M tok/cr",
            "snowflake-arctic \u2014 1.2M tok/cr",
        ], "Complex reasoning, nuanced analysis"),
        ("Premium", DANGER, [
            "claude-3-5-sonnet \u2014 667K/133K",
            "llama3.1-405b \u2014 333K tok/cr",
            "mistral-large \u2014 196K tok/cr",
        ], "Highest quality, coding, long-form"),
    ]

    for i, (tier_name, color, models, desc) in enumerate(tiers):
        left = left_start + i * (col_w + gap)
        shape = add_card(slide, left, top, col_w, col_h,
                         tier_name, models, title_color=color,
                         border_color=color, border_width=2)
        # Add description at bottom of card
        add_textbox(slide, left + Inches(0.1), top + col_h - Inches(0.45),
                    col_w - Inches(0.2), Inches(0.35),
                    desc, font_size=9, color=color, alignment=PP_ALIGN.LEFT)

    add_highlight_box(slide,
                      "The right model is the cheapest one that meets your quality bar.",
                      Inches(0.7), Inches(5.0), Inches(11.5))

    add_footer(slide)


def build_slide_06_token_economics(prs):
    """SLIDE 6: Token Economics & Shadow Waste."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    add_slide_title(slide, "Token Economics & Shadow Waste")
    add_subtitle(slide, "Module 01 Sections C-D \u2014 From tokens to dollars, and where savings hide")

    # Left card — Token Cost Estimator
    add_card(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(3.0),
             "Token Cost Estimator",
             ["Estimate daily/monthly credit burn per model",
              "Formula: rows \u00d7 avg_tokens \u00f7 tokens_per_credit",
              "Dollar projection at $3/credit",
              "Same 10K rows: mistral-7b = 0.001 cr vs claude-3-5-sonnet = 0.075 cr",
              "Interactive calculator in Module 01 Section D"],
             title_color=BLUE)

    # Right — Shadow Waste Savings table
    savings_data = [
        ["Task", "COMPLETE()", "Specialized", "Savings"],
        ["Sentiment", "0.30 cr", "0.02 cr", "93%"],
        ["Translate", "0.45 cr", "0.05 cr", "89%"],
        ["Classify", "0.30 cr", "0.03 cr", "90%"],
        ["Extract", "0.30 cr", "0.04 cr", "87%"],
    ]

    add_textbox(slide, Inches(6.7), Inches(1.6), Inches(5.5), Inches(0.35),
                "Shadow Waste Savings", font_size=13, bold=True, color=BLUE)

    add_table(slide, savings_data,
              Inches(6.7), Inches(2.1), Inches(5.5), Inches(2.2),
              col_widths=[Inches(1.3), Inches(1.4), Inches(1.4), Inches(1.4)])

    # Warning about COMPLETE() shadow waste
    add_textbox(slide, Inches(6.7), Inches(4.4), Inches(5.5), Inches(0.4),
                "Using COMPLETE() for tasks with specialized functions is the #1 source of shadow waste.",
                font_size=9, color=ORANGE)

    # Bottom highlight
    add_highlight_box(slide,
                      "Shadow Waste = AI credit consumption that delivers little or no business value. "
                      "8 detection patterns in Module 06.",
                      Inches(0.7), Inches(5.1), Inches(11.5))

    add_footer(slide)


def build_slide_07_setup(prs):
    """SLIDE 7: Environment Setup — Modules 03-04."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    add_slide_title(slide, "Environment Setup")
    add_subtitle(slide, "Modules 03-04 \u2014 Cortex Code CLI + Governed Lab Environment")

    # Left card — CLI
    add_card(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(2.8),
             "Cortex Code CLI (Module 03)",
             ["Install via curl or PowerShell",
              "Connect to Snowflake account",
              "First prompt: natural language \u2192 SQL",
              "Commands: /status, /help, #TABLE"],
             title_color=BLUE)

    # Right card — Lab
    add_card(slide, Inches(6.7), Inches(1.6), Inches(5.5), Inches(2.8),
             "Lab Environment (Module 04)",
             ["Database: cortex_lab \u00b7 Schema: ai_workshop",
              "Warehouse: cortex_wh (SMALL, 60s suspend)",
              "Resource Monitor: 50 credits, 90% suspend",
              "Role: cortex_analyst",
              "Sample: 500-row customer_feedback"],
             title_color=BLUE)

    # Warning
    add_warning_box(slide,
                    "Resource Monitor guards warehouse compute credits, not AI token credits. "
                    "AI credits are billed separately \u2014 use Budgets and "
                    "METERING_DAILY_HISTORY to track AI spend.",
                    Inches(0.7), Inches(4.8), Inches(11.5), Inches(0.65))

    add_footer(slide)


def build_slide_08_exercises(prs):
    """SLIDE 8: AI SQL Hands-On Exercises — Module 05."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    add_slide_title(slide, "AI SQL Hands-On Exercises")
    add_subtitle(slide, "Module 05 \u2014 20 min \u2014 Run Cortex functions and compare costs")

    exercises = [
        (Inches(0.7), Inches(1.6), "Ex 1: AI_SENTIMENT()",
         ["Analyze customer feedback sentiment", "No model param needed"],
         BLUE, "Lowest Cost", GREEN),
        (Inches(6.7), Inches(1.6), "Ex 2: AI_CLASSIFY()",
         ["Route feedback to support teams", "Zero-shot classification"],
         BLUE, "Lower Cost", GREEN),
        (Inches(0.7), Inches(3.4), "Ex 3: AI_COMPLETE()",
         ["Extract structured info from text", "Using budget model mistral-7b"],
         BLUE, "mistral-7b", BLUE),
        (Inches(6.7), Inches(3.4), "Ex 4: Model Cost Comparison",
         ["Same prompt, 3 models:", "mistral-7b vs llama3-70b vs claude-3-5-sonnet"],
         DANGER, None, None),
    ]

    for left, top, title, bullets, title_clr, badge, badge_clr in exercises:
        border_clr = DANGER if title_clr == DANGER else CARD_BORDER
        bw = 2 if title_clr == DANGER else 1
        add_card(slide, left, top, Inches(5.5), Inches(1.5),
                 title, bullets, title_color=title_clr,
                 border_color=border_clr, border_width=bw)

    # Ex 5 — full width
    add_card(slide, Inches(0.7), Inches(5.2), Inches(11.5), Inches(1.0),
             "Ex 5: Scale Simulation",
             ["Test on 5 rows \u2192 measure cost \u2192 project to 500K rows \u2192 decide"],
             title_color=BLUE)

    add_footer(slide)


def build_slide_09_cost_tracking(prs):
    """SLIDE 9: Token & Credit Usage Tracking — Module 06."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    add_slide_title(slide, "Token & Credit Usage Tracking")
    add_subtitle(slide, "Module 06 \u2014 13 queries \u2014 Your real-time FinOps console")

    queries = [
        (Inches(0.7), Inches(1.6), "Q1\u2013Q5: Core Tracking",
         ["AI credits by function, model, hour", "Cost projection from dedicated views"],
         BLUE, CARD_BORDER),
        (Inches(6.7), Inches(1.6), "Q6: Idle Search Detection",
         ["Search indexing with zero queries", "100% waste (Transition \u2192 Module 07)"],
         BLUE, CARD_BORDER),
        (Inches(0.7), Inches(3.0), "Q7: Over-Sized Models",
         ["credits_per_1M_tokens metric", "Find where cheaper models suffice (40\u201375% savings)"],
         BLUE, DANGER),
        (Inches(6.7), Inches(3.0), "Q8\u2013Q9: Redundant Calls & Prompt Bloat",
         ["Duplicate pipelines detected", "Rising avg tokens/call over time"],
         BLUE, DANGER),
        (Inches(0.7), Inches(4.4), "Q10\u2013Q11: Agent Loops & Dev in Prod",
         ["MAX vs AVG credit gap", "LATERAL FLATTEN(ROLE_NAMES)"],
         BLUE, DANGER),
        (Inches(6.7), Inches(4.4), "Q12\u2013Q13: Anomaly & Attribution",
         ["WoW LAG(7) spike detection", "Tagged vs untagged spend ratio"],
         BLUE, DANGER),
    ]

    for left, top, title, bullets, title_clr, border_clr in queries:
        add_card(slide, left, top, Inches(5.5), Inches(1.2),
                 title, bullets, title_color=title_clr,
                 border_color=border_clr, border_width=2 if border_clr == DANGER else 1)

    add_warning_box(slide,
                    "Shadow Waste: Q7\u2013Q13 detect 8 patterns of hidden AI credit drain. "
                    "AI services are serverless \u2014 no warehouse to watch, no query fails. "
                    "Credits silently accumulate.",
                    Inches(0.7), Inches(5.9), Inches(11.5), Inches(0.6))

    add_footer(slide)


def build_slide_10_ai_credits(prs):
    """SLIDE 10: AI Credits Transition — Module 07."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    add_slide_title(slide, "AI Credits Transition")
    add_subtitle(slide, "Module 07 \u2014 5 min \u2014 New SERVICE_TYPE billing you need to track")

    # Left — SERVICE_TYPE table
    st_data = [
        ["SERVICE_TYPE", "Description"],
        ["CORTEX_AGENTS", "Agentic AI workflows"],
        ["CORTEX_CODE_CLI", "Cortex Code CLI usage"],
        ["CORTEX_CODE_SNOWSIGHT", "Cortex Code in browser"],
        ["SNOWFLAKE_INTELLIGENCE", "Snowflake Intelligence"],
    ]

    add_textbox(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.35),
                "4 New SERVICE_TYPEs", font_size=13, bold=True, color=ORANGE)

    add_table(slide, st_data,
              Inches(0.7), Inches(2.1), Inches(5.5), Inches(2.5),
              col_widths=[Inches(3.0), Inches(2.5)])

    # Right — Detection Query code block
    add_textbox(slide, Inches(6.7), Inches(1.6), Inches(5.5), Inches(0.35),
                "Detection Query", font_size=13, bold=True, color=BLUE)

    sql_code = """SELECT SERVICE_TYPE,
       SUM(CREDITS_USED) AS credits,
       credits * 3 AS est_cost
FROM SNOWFLAKE.ACCOUNT_USAGE
    .METERING_DAILY_HISTORY
WHERE SERVICE_TYPE LIKE 'CORTEX%'
   OR SERVICE_TYPE = 'SNOWFLAKE_INTELLIGENCE'
GROUP BY 1 ORDER BY 2 DESC;"""

    add_code_box(slide, sql_code,
                 Inches(6.7), Inches(2.1), Inches(5.5), Inches(2.5))

    # Warning
    add_warning_box(slide,
                    "These credits appear in METERING_DAILY_HISTORY but NOT in "
                    "CORTEX_AI_FUNCTIONS_USAGE_HISTORY. You must check both views "
                    "for complete AI cost visibility.",
                    Inches(0.7), Inches(4.8), Inches(11.5), Inches(0.5))

    # Credit Guardrails highlight
    guardrails_data = [
        ["Layer", "Tool", "Behavior"],
        ["Proactive", "Per-user credit limits", "Hard block \u2014 rolling 24h"],
        ["Threshold", "Monthly Budgets", "Alert before overshoot"],
        ["Audit", "Usage History views", "Attribution by user + model"],
    ]

    add_textbox(slide, Inches(0.7), Inches(5.5), Inches(5.5), Inches(0.3),
                "Credit Guardrails", font_size=11, bold=True, color=ORANGE)

    add_table(slide, guardrails_data,
              Inches(0.7), Inches(5.85), Inches(5.5), Inches(1.2),
              col_widths=[Inches(1.3), Inches(2.2), Inches(2.0)])

    add_textbox(slide, Inches(6.7), Inches(5.5), Inches(5.5), Inches(0.7),
                "Parameters:\n"
                "CORTEX_CODE_CLI_DAILY_EST_CREDIT_LIMIT_PER_USER\n"
                "CORTEX_CODE_SNOWSIGHT_DAILY_EST_CREDIT_LIMIT_PER_USER",
                font_size=9, color=MUTED)

    add_footer(slide)


def build_slide_11_dashboard(prs):
    """SLIDE 11: Streamlit FinOps Dashboard — Module 08."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    add_slide_title(slide, "Streamlit FinOps Dashboard")
    add_subtitle(slide, "Module 08 \u2014 13 min \u2014 Build with Cortex Code, deploy to Snowflake")

    cards = [
        (Inches(0.7), Inches(1.6), "KPI + Visualizations",
         ["Total Queries, Credits, Avg/Call, Users, Est. $",
          "Hourly Trend \u00b7 Function Breakdown",
          "Model Cost Comparison \u00b7 Top Expensive Queries"],
         BLUE, CARD_BORDER),
        (Inches(6.7), Inches(1.6), "Attribution Coverage",
         ["Tagged vs Untagged donut chart",
          "% coverage metric (target: >90%)"],
         DANGER, DANGER),
        (Inches(0.7), Inches(3.7), "WoW Anomaly Detection",
         ["Daily spend vs same day last week",
          "Color-coded: >50% yellow, >200% red"],
         DANGER, DANGER),
        (Inches(6.7), Inches(3.7), "Sample Dashboard Output",
         ["CSS mockup preview in Module 08 Section C",
          "5 KPI cards \u00b7 Donut \u00b7 Area chart \u00b7 Cost grid"],
         BLUE, BLUE),
    ]

    for left, top, title, bullets, title_clr, border_clr in cards:
        add_card(slide, left, top, Inches(5.5), Inches(1.8),
                 title, bullets, title_color=title_clr,
                 border_color=border_clr, border_width=2)

    add_highlight_box(slide,
                      "Deploy: Snowsight \u2192 Projects \u2192 Streamlit \u2192 + Streamlit App "
                      "\u2192 Paste code \u2192 Run",
                      Inches(0.7), Inches(5.8), Inches(11.5))

    add_footer(slide)


def build_slide_12_quick_ref(prs):
    """SLIDE 12: Quick Reference — 7 ACCOUNT_USAGE Views."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    add_slide_title(slide, "Quick Reference: 7 ACCOUNT_USAGE Views")
    add_subtitle(slide, "Module 09 Section E \u2014 The views you need for complete AI cost visibility")

    ref_data = [
        ["View", "Tracks", "Key Columns"],
        ["CORTEX_AI_FUNCTIONS_USAGE_HISTORY", "LLM / task functions",
         "FUNCTION_NAME, MODEL_NAME, CREDITS, USER_ID"],
        ["CORTEX_AISQL_USAGE_HISTORY", "AI SQL tokens",
         "TOKENS, TOKEN_CREDITS, QUERY_TAG, USER_ID"],
        ["DOCUMENT_AI_USAGE_HISTORY", "Document AI",
         "CREDITS_USED, PAGE_COUNT"],
        ["CORTEX_SEARCH_DAILY_USAGE_HISTORY", "Search services",
         "SERVICE_NAME, CREDITS"],
        ["METERING_DAILY_HISTORY", "All SERVICE_TYPEs",
         "SERVICE_TYPE, CREDITS_BILLED"],
        ["CORTEX_CODE_CLI_USAGE_HISTORY", "Cortex Code CLI",
         "CREDITS, MODEL_NAME, USER_ID"],
        ["CORTEX_CODE_SNOWSIGHT_USAGE_HISTORY", "Cortex Code Snowsight",
         "CREDITS, MODEL_NAME, USER_ID"],
    ]

    add_table(slide, ref_data,
              Inches(0.7), Inches(1.6), Inches(11.5), Inches(3.5),
              col_widths=[Inches(5.0), Inches(2.2), Inches(4.3)])

    add_warning_box(slide,
                    "USER_ID in Cortex views is a NUMBER. JOIN to "
                    "SNOWFLAKE.ACCOUNT_USAGE.USERS for names. ~45 min latency on all views.",
                    Inches(0.7), Inches(5.4), Inches(11.5), Inches(0.55))

    add_footer(slide)


def build_slide_13_finops_habits(prs):
    """SLIDE 13: Three FinOps Habits + AI Governance Principles."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    add_slide_title(slide, "The Three FinOps Habits for AI")

    habits = [
        ("ESTIMATE", BLUE, "Test on 10 rows. Project to full dataset. Decide."),
        ("MONITOR", GREEN, "Dashboard is your early warning system. Check weekly."),
        ("CHOOSE", ORANGE, "The cheapest model that meets your quality bar."),
    ]

    habit_w = Inches(3.7)
    for i, (name, color, desc) in enumerate(habits):
        left = Inches(0.7) + i * (habit_w + Inches(0.2))
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), habit_w, Inches(1.4)
        )
        set_shape_fill(shape, GRAY_BG)
        set_shape_border(shape, color, 2)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = FONT_HEADING
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(10)
        run2.font.color.rgb = TEXT_COLOR
        run2.font.name = FONT_BODY

    # AI Governance Principles
    add_textbox(slide, Inches(0.7), Inches(3.0), Inches(11.5), Inches(0.35),
                "AI Governance Principles (Module 09 Section C)",
                font_size=14, bold=True, color=BLUE)

    principles = [
        ("Credit Burn", DANGER,
         "AI costs scale with data volume, not warehouse size"),
        ("Behavior", ORANGE,
         "Model choice & prompt design drive cost, not infra"),
        ("Visibility", BLUE,
         "Dedicated views, not query_history, for AI tracking"),
        ("Signals", GREEN,
         "Shadow waste hides where traditional monitors can\u2019t see"),
    ]

    card_w = Inches(2.7)
    for i, (name, color, desc) in enumerate(principles):
        left = Inches(0.7) + i * (card_w + Inches(0.2))
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.5), card_w, Inches(1.8)
        )
        set_shape_fill(shape, GRAY_BG)
        set_shape_border(shape, CARD_BORDER, 1)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = FONT_HEADING
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(9)
        run2.font.color.rgb = TEXT_COLOR
        run2.font.name = FONT_BODY

    add_footer(slide)


def build_slide_14_closing(prs):
    """SLIDE 14: What You Built Today — blue gradient, pipeline grid, closing quote."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # BLANK

    # Solid blue background (approximates gradient)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLUE

    # Title
    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.6),
                "What You Built Today", font_size=28, bold=True,
                color=WHITE, alignment=PP_ALIGN.CENTER)

    # Pipeline grid — 3 cols x 3 rows (7 items)
    modules = [
        ("01-02", "Token Economy & Cortex AI"),
        ("03-04", "Environment & Setup"),
        ("05", "AI SQL Hands-On"),
        ("06", "Cost Tracking"),
        ("07", "AI Credits Transition"),
        ("08", "FinOps Dashboard"),
        ("09", "Govern & Scale"),
    ]

    cell_w = Inches(3.5)
    cell_h = Inches(1.0)
    gap = Inches(0.25)
    grid_left = Inches(1.5)
    grid_top = Inches(1.3)

    for idx, (mod_num, mod_name) in enumerate(modules):
        col = idx % 3
        row = idx // 3
        left = grid_left + col * (cell_w + gap)
        top = grid_top + row * (cell_h + gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, cell_w, cell_h
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x06, 0x9E, 0xD8)  # slightly lighter blue
        shape.fill.fore_color.rgb = RGBColor(0x06, 0x9E, 0xD8)
        # Make slightly translucent look via lighter blue
        set_shape_border(shape, WHITE, 1)
        shape.shadow.inherit = False

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = mod_num
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xB0, 0xE0, 0xFF)
        run.font.name = FONT_BODY
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = mod_name
        run2.font.size = Pt(12)
        run2.font.bold = True
        run2.font.color.rgb = WHITE
        run2.font.name = FONT_BODY

    # Closing quote
    add_textbox(slide, Inches(1.0), Inches(4.5), Inches(11.0), Inches(0.6),
                "AI spend will grow. The teams who govern it now will be the "
                "ones who can afford to scale it later.",
                font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Logo
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(5.9), Inches(5.3), Inches(1.4))


# ════════════════════════════════════════════════════════════════
#  Main
# ════════════════════════════════════════════════════════════════

def main():
    print("Loading template:", TEMPLATE)
    prs = Presentation(TEMPLATE)

    print("Removing existing template slides...")
    delete_existing_slides(prs)

    print("Building 14 slides...")
    build_slide_01_title(prs)
    print("  Slide  1: Title")
    build_slide_02_agenda(prs)
    print("  Slide  2: Agenda")
    build_slide_03_cost_layers(prs)
    print("  Slide  3: Two Cost Layers")
    build_slide_04_capability_map(prs)
    print("  Slide  4: Capability Map")
    build_slide_05_model_tiers(prs)
    print("  Slide  5: Model Cost Tiers")
    build_slide_06_token_economics(prs)
    print("  Slide  6: Token Economics")
    build_slide_07_setup(prs)
    print("  Slide  7: Environment Setup")
    build_slide_08_exercises(prs)
    print("  Slide  8: Hands-On Exercises")
    build_slide_09_cost_tracking(prs)
    print("  Slide  9: Cost Tracking")
    build_slide_10_ai_credits(prs)
    print("  Slide 10: AI Credits Transition")
    build_slide_11_dashboard(prs)
    print("  Slide 11: Streamlit Dashboard")
    build_slide_12_quick_ref(prs)
    print("  Slide 12: Quick Reference (7 Views)")
    build_slide_13_finops_habits(prs)
    print("  Slide 13: FinOps Habits + Governance")
    build_slide_14_closing(prs)
    print("  Slide 14: What You Built Today")

    print(f"\nSaving to: {OUTPUT}")
    prs.save(OUTPUT)
    print(f"Done! {len(prs.slides)} slides written.")


if __name__ == "__main__":
    main()
